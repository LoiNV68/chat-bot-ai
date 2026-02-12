from typing import List, Dict, Any
from app.services.optimized_excel_parser import OptimizedExcelParser
from fastapi import UploadFile
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, update
from app.models.document import Document, AccessScope
from app.services.vector_store import VectorStore
from app.schemas.doc_schema import DocUploadSchema
import pandas as pd
import pdfplumber
import uuid
import datetime
import datetime
import os

# Cache toàn cục cho EasyOCR reader
ocr_reader = None

# Tối ưu bộ nhớ PyTorch cho môi trường GPU dùng chung
os.environ["PYTORCH_CUDA_ALLOC_CONF"] = "expandable_segments:True"

def get_ocr_reader():
    global ocr_reader
    if ocr_reader is None:
        import easyocr
        import torch
        
        use_gpu = torch.cuda.is_available()
        gpu_name = torch.cuda.get_device_name(0) if use_gpu else "CPU"
        print(f"[DEBUG] Initializing EasyOCR Reader (vi, en)... GPU Detected: {use_gpu} ({gpu_name})")
        
        ocr_reader = easyocr.Reader(['vi', 'en'], gpu=use_gpu)
    return ocr_reader

def perform_ocr_on_image(image_obj):
    """
    Hàm hỗ trợ để chạy trong thread riêng.
    Nhận vào PIL Image hoặc bytes và trả về text.
    """
    try:
        reader = get_ocr_reader()
        import numpy as np
        # Chuyển đổi PIL Image thành numpy array
        result = reader.readtext(np.array(image_obj), detail=0)
        return "\n".join(result)
    except Exception as e:
        print(f"[DEBUG] OCR Error on GPU: {e}")
        if "out of memory" in str(e).lower():
             print("[DEBUG] Falling back to CPU for OCR due to OOM...")
             torch.cuda.empty_cache()
             try:
                 import easyocr
                  # Khởi tạo reader mới trên CPU
                 cpu_reader = easyocr.Reader(['vi', 'en'], gpu=False)
                 result = cpu_reader.readtext(np.array(image_obj), detail=0)
                 return "\n".join(result)
             except Exception as cpu_e:
                 print(f"[DEBUG] OCR CPU Fallback Error: {cpu_e}")
        return ""

class IngestionService:
    def __init__(self, db: AsyncSession):
        self.db = db
        self.vector_store = VectorStore()
        self.excel_parser = OptimizedExcelParser(rows_per_chunk=3, overlap=1)

    async def save_document(self, file: UploadFile, metadata: DocUploadSchema, user_id: int):
        # 1. Đọc Nội dung File Trước
        try:
            content = await file.read()
            if not content:
                raise ValueError("Empty file content")
        except Exception as e:
            return {"status": "error", "message": f"Failed to read file: {str(e)}"}

        # 2. Lưu file vào đĩa
        from datetime import datetime
        
        # Sử dụng đường dẫn tuyệt đối tương đối với thư mục gốc project (giả sử chạy từ thư mục gốc project)
        base_dir = os.getcwd()
        upload_dir = os.path.join(base_dir, "uploads")
        os.makedirs(upload_dir, exist_ok=True)
        
        # Sử dụng timestamp để tránh xung đột tên
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        safe_filename = f"{timestamp}_{file.filename}"
        file_path = os.path.join(upload_dir, safe_filename)
        
        try:
            with open(file_path, "wb") as f:
                f.write(content)
        except Exception as e:
             return {"status": "error", "message": f"Failed to save file: {str(e)}"}

        # 3. Kiểm tra Version Control & Giao dịch DB
        try:
            existing_doc_query = select(Document).where(
                Document.filename == file.filename,
                Document.is_active == True
            )
            result = await self.db.execute(existing_doc_query)
            existing_doc = result.scalars().first()
            
            new_version = 1
            parent_id = None
            
            if existing_doc:
                existing_doc.is_active = False
                new_version = existing_doc.version + 1
                parent_id = existing_doc.id
                # Chỉ đánh dấu inactive ở đây, commit sau
                
            # 4. Tạo Bản ghi Tài liệu Mới
            new_doc = Document(
                filename=file.filename,
                version=new_version,
                is_active=True,
                effective_date=metadata.effective_date,
                expiry_date=metadata.expiry_date,
                parent_id=parent_id,
                access_scope=metadata.scope,
                target_id=metadata.target_id,
                uploaded_by=user_id,
                file_path=file_path
            )
            self.db.add(new_doc)
            await self.db.commit() # Commit giao dịch DB
            await self.db.refresh(new_doc)
            
            return {
                "status": "success", 
                "doc_id": new_doc.id, 
                "file_path": file_path,
                "filename": file.filename,
                "version": new_version,
                "scope": metadata.scope.value,
                "target_id": metadata.target_id,
                "message": "Upload successful. AI processing in background."
            }
            
        except Exception as e:
            await self.db.rollback()
            # Dọn dẹp file
            if os.path.exists(file_path):
                os.remove(file_path)
            return {"status": "error", "message": f"Database error: {str(e)}"}



    async def ingest_file_content(self, doc_id: int, file_path: str, filename: str, version: int, scope: str, target_id: str):
        print(f"[DEBUG] Starting ingestion for doc_id={doc_id}, file={filename}")
        # Đọc nội dung từ ĐĨA
        try:
            if not os.path.exists(file_path):
                 print(f"[DEBUG] File NOT FOUND at {file_path}")
                 return

            with open(file_path, "rb") as f:
                content = f.read()
            print(f"[DEBUG] Read {len(content)} bytes from {file_path}")
        except Exception as e:
            print(f"[DEBUG] Background ingestion failed to read file {file_path}: {e}")
            return

        # Phân tích Nội dung
        chunks = []
        ext = filename.lower()
        
        # Danh sách các định dạng được hỗ trợ
        SUPPORTED_EXTENSIONS = {
            'excel': ['.xlsx', '.xls'],
            'pdf': ['.pdf'],
            'word': ['.docx', '.doc'],
            'text': ['.txt', '.md', '.rtf'],
            'csv': ['.csv'],
            'powerpoint': ['.pptx', '.ppt']
        }
        
        try:
            # Excel files (.xlsx, .xls)
            if any(ext.endswith(e) for e in SUPPORTED_EXTENSIONS['excel']):
                print(f"[DEBUG] Parsing Excel file: {filename}")
                chunks = self._parse_excel(content, filename)
            
            # PDF files
            elif ext.endswith('.pdf'):
                print("[DEBUG] Parsing PDF (checking for OCR)...")
                chunks = await self._parse_pdf(content, filename)
            
            # Word files (.docx, .doc)
            elif any(ext.endswith(e) for e in SUPPORTED_EXTENSIONS['word']):
                print(f"[DEBUG] Parsing Word file: {filename}")
                chunks = self._parse_docx(content, filename)
            
            # Text files (.txt, .md, .rtf)
            elif any(ext.endswith(e) for e in SUPPORTED_EXTENSIONS['text']):
                print(f"[DEBUG] Parsing Text file: {filename}")
                chunks = self._parse_text(content, filename)
            
            # CSV files
            elif ext.endswith('.csv'):
                print(f"[DEBUG] Parsing CSV file: {filename}")
                chunks = self._parse_csv(content, filename)
            
            # PowerPoint files (.pptx, .ppt)
            elif any(ext.endswith(e) for e in SUPPORTED_EXTENSIONS['powerpoint']):
                print(f"[DEBUG] Parsing PowerPoint file: {filename}")
                chunks = self._parse_pptx(content, filename)
            
            else:
                all_supported = [e for exts in SUPPORTED_EXTENSIONS.values() for e in exts]
                print(f"[DEBUG] Skipping Qdrant ingestion for unsupported file type: {filename}")
                print(f"[DEBUG] Supported formats: {', '.join(all_supported)}")
                return
            
            print(f"[DEBUG] Extracted {len(chunks)} chunks.")
        except Exception as e:
             # Bắt full traceback để debug
             import traceback
             traceback.print_exc()
             print(f"[DEBUG] Parsing error for {filename}: {e}")
             return

        if not chunks:
             print(f"[DEBUG] Warning: No text content found in {filename}")
             return

        texts = [chunk["content"] for chunk in chunks]
        metadatas_list = [chunk["metadata"] for chunk in chunks]
        
        # Cập nhật metadata mà không cần truy vấn DB nếu có thể, hoặc giả sử dữ liệu truyền vào đúng
        for i, meta in enumerate(metadatas_list):
            meta.update({
                "doc_id": doc_id,
                "version": version,
                "access_scope": scope,
                "target_id": target_id,
                "content": texts[i],
                "is_active": True # QUAN TRỌNG cho ChatEngine filtering 
            })
            
        ids = [str(uuid.uuid4()) for _ in range(len(texts))]
        
        print(f"[DEBUG] Upserting {len(texts)} vectors to Qdrant...")
        try:
            await self.vector_store.upsert_vectors(texts, metadatas_list, ids)
            print(f"[DEBUG] Background ingestion COMPLETED for doc {doc_id}.")
            
            # Cập nhật cờ is_processed trong database
            async with self.db.begin():
                await self.db.execute(
                    update(Document).where(Document.id == doc_id).values(is_processed=True)
                )
            print(f"[DEBUG] Document {doc_id} marked as processed.")
        except Exception as e:
            print(f"[DEBUG] Vector upsert failed: {e}")

    def _extract_time_from_text(self, text: str) -> str:
        """
        Trích xuất thông tin Học kỳ / Năm học từ văn bản.
        Hỗ trợ:
        - Năm học 20xx-20yy
        - Năm học 20xx - 20yy
        - Năm 20xx - 20yy
        - 20xx - 20yy (nếu có từ khóa Học kỳ)
        - Năm 25-26 (viết tắt)
        """
        import re
        text = text.replace('\n', ' ')
        
        # Helper để chuẩn hóa năm 2 số -> 4 số
        def normalize_year(y):
            if len(y) == 2:
                return "20" + y
            return y

        # 1. Tìm pattern đầy đủ: "Học kỳ ... Năm học ...-..."
        # Pattern: Học kỳ [I,II,1,2] ... Năm (học)? 20xx(-|_)20yy
        full_pattern = re.search(r'(?:Học kỳ|Kỳ|HK)\s*([0-9IVX]+).*?(?:Năm học|Năm)\s*(\d{2,4})[_\-](\d{2,4})', text, re.IGNORECASE)
        if full_pattern:
            hk_str = full_pattern.group(1).upper()
            y1 = normalize_year(full_pattern.group(2))
            y2 = normalize_year(full_pattern.group(3))
            
            roman_map = {'I': '1', 'II': '2', 'III': '3', 'IV': '4'}
            hk = roman_map.get(hk_str, hk_str)
            
            return f"**Học kỳ {hk}, Năm học {y1}-{y2}**"
            
        # 2. Tìm riêng lẻ
        # Tìm Học kỳ
        hk_match = re.search(r'(?:Học kỳ|Kỳ|HK)\s*([0-9IVX]+)', text, re.IGNORECASE)
        hk_info = ""
        if hk_match:
            hk_str = hk_match.group(1).upper()
            roman_map = {'I': '1', 'II': '2', 'III': '3', 'IV': '4'}
            hk = roman_map.get(hk_str, hk_str)
            hk_info = f"Học kỳ {hk}"
            
        # Tìm Năm học (linh hoạt hơn)
        # Hỗ trợ: "Năm học 2024-2025", "Năm 24-25", "NH 24-25"
        year_match = re.search(r'(?:Năm học|Năm|NH)\s*(\d{2,4})[_\-](\d{2,4})', text, re.IGNORECASE)
        year_info = ""
        if year_match:
            y1 = normalize_year(year_match.group(1))
            y2 = normalize_year(year_match.group(2))
            year_info = f"Năm học {y1}-{y2}"
        
        # Kết hợp
        if hk_info and year_info:
            return f"**{hk_info}, {year_info}**"
        elif year_info:
            return f"**{year_info}**"
        elif hk_info:
            return f"**{hk_info}**"
            
        return ""

    def _parse_excel(self, content: bytes, filename: str) -> List[Dict]:
        """Parse Excel files sử dụng OptimizedExcelParser (4 chiến lược chunking)"""
        return self.excel_parser.parse_excel(content, filename)


    async def _parse_pdf(self, content: bytes, filename: str) -> List[Dict]:
        import pdfplumber
        import io
        import asyncio
        
        chunks = []
        extracted_time = ""
        
        try:
            # Phải tải bytes vào IO. 
            # Lưu ý: pdfplumber.open là sync. Nếu file rất lớn, có thể block một chút, 
            # nhưng thường thì parsing là bottleneck.
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                # --- EXTRACT TIME INFO FROM FIRST PAGE ---
                if len(pdf.pages) > 0:
                    first_page_text = pdf.pages[0].extract_text() or ""
                    extracted_time = self._extract_time_from_text(first_page_text)
                    print(f"[DEBUG] Extracted time from PDF content ({filename}): {extracted_time}")
                # -----------------------------------------

                for page in pdf.pages:
                    # 1. Thử Trích xuất Text Chuẩn
                    text = page.extract_text() or ""
                    
                    if text.strip():
                         metadata = {"source": filename, "page": page.page_number, "type": "text"}
                         if extracted_time:
                             metadata["time_info"] = extracted_time
                         
                         chunks.append({
                             "content": text,
                             "metadata": metadata
                         })
                    
                    # 2. Trích xuất Bảng
                    tables = page.extract_tables()
                    for table in tables:
                        table_str = str(table) 
                        
                        metadata = {"source": filename, "page": page.page_number, "type": "table"}
                        if extracted_time:
                             metadata["time_info"] = extracted_time

                        chunks.append({
                             "content": table_str,
                             "metadata": metadata
                         })
                    
                    # 3. Fallback: Logic OCR nếu text trống/ít
                    if len(text.strip()) < 50:
                        # Chỉ thử render và OCR nếu thực sự cần.
                        print(f"[DEBUG] Page {page.page_number} has low text. Attempting OCR...")
                        
                        try:
                             # Render page thành image. Có thể chậm/nặng, có thể offload?
                             # page.to_image() with resolution=300 is CPU intensive.
                             # Let's run this in thread too if possible, but page object pickling is tricky.
                             # We'll run to_image in main loop (it uses C extensions often, release GIL?)
                             # But let's try to trust it's fast enough or offload the WHOLE thing?
                             # Hiện tại, chỉ offload phần OCR vì nó nặng nhất.
                             
                             # Cần bắt sớm nếu thiếu dependencies
                             try:
                                 # resolution=200 nhanh hơn 300 và thường đủ
                                 # Sử dụng context manager cho to_image đôi khi giúp cleanup?
                                 p_im = page.to_image(resolution=200)
                                 original_image = p_im.original
                             except Exception as render_err:
                                 print(f"[DEBUG] Page rendering failed (missing valid backend?): {render_err}")
                                 continue
                             
                             # Chạy OCR trong thread riêng
                             ocr_text = await asyncio.to_thread(perform_ocr_on_image, original_image)
                             
                             if ocr_text.strip():
                                 print(f"[DEBUG] OCR Success on page {page.page_number}: extracted {len(ocr_text)} chars.")
                                 
                                 metadata = {"source": filename, "page": page.page_number, "type": "ocr_text"}
                                 # If extracted_time was not found in standard text, try finding in OCR text of first page
                                 if not extracted_time and page.page_number == 1:
                                     extracted_time = self._extract_time_from_text(ocr_text)
                                     print(f"[DEBUG] Extracted time from PDF OCR content ({filename}): {extracted_time}")
                                 
                                 if extracted_time:
                                     metadata["time_info"] = extracted_time

                                 chunks.append({
                                     "content": ocr_text,
                                     "metadata": metadata
                                 })
                        except Exception as e:
                            print(f"[DEBUG] OCR process failed on page {page.page_number}: {e}")

        except Exception as e:
            print(f"[DEBUG] Error parsing PDF: {e}")
            import traceback
            traceback.print_exc()
            
        return chunks

    def _parse_docx(self, content: bytes, filename: str) -> List[Dict]:
        import io
        from docx import Document as DocxDocument
        
        chunks = []
        try:
            doc = DocxDocument(io.BytesIO(content))
            # Trích xuất các đoạn văn
            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    chunks.append({
                        "content": text,
                        "metadata": {"source": filename, "type": "text", "paragraph": i}
                    })
            
            # Trích xuất các bảng
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                
                table_str = "\n".join(table_data)
                if table_str.strip():
                     chunks.append({
                        "content": table_str,
                        "metadata": {"source": filename, "type": "table", "table_index": i}
                    })
                    
        except Exception as e:
            print(f"[DEBUG] Error parsing DOCX: {e}")
        
        return chunks

    def _parse_text(self, content: bytes, filename: str) -> List[Dict]:
        """Parse text files (.txt, .md, .rtf)"""
        chunks = []
        try:
            # Try different encodings
            text = None
            for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                try:
                    text = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                print(f"[DEBUG] Could not decode text file {filename}")
                return chunks
            
            # Split by paragraphs (double newline) or by reasonable chunk size
            paragraphs = text.split('\n\n')
            for i, para in enumerate(paragraphs):
                para = para.strip()
                if para:
                    chunks.append({
                        "content": para,
                        "metadata": {"source": filename, "type": "text", "paragraph": i}
                    })
        except Exception as e:
            print(f"[DEBUG] Error parsing text file: {e}")
        
        return chunks

    def _parse_csv(self, content: bytes, filename: str) -> List[Dict]:
        """Parse CSV files"""
        import io
        chunks = []
        try:
            # Try different encodings
            for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                try:
                    df = pd.read_csv(io.BytesIO(content), encoding=encoding)
                    break
                except (UnicodeDecodeError, pd.errors.ParserError):
                    continue
            else:
                print(f"[DEBUG] Could not parse CSV file {filename}")
                return chunks
            
            for _, row in df.iterrows():
                row_str = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                if row_str.strip():
                    chunks.append({
                        "content": row_str,
                        "metadata": {"source": filename, "type": "csv"}
                    })
        except Exception as e:
            print(f"[DEBUG] Error parsing CSV: {e}")
        
        return chunks

    def _parse_pptx(self, content: bytes, filename: str) -> List[Dict]:
        """Parse PowerPoint files (.pptx, .ppt)"""
        import io
        chunks = []
        
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    chunks.append({
                        "content": "\n".join(slide_text),
                        "metadata": {"source": filename, "type": "powerpoint", "slide": slide_num}
                    })
        except ImportError:
            print("[DEBUG] python-pptx not installed. Install with: pip install python-pptx")
        except Exception as e:
            # Fallback for .ppt (old format) - might need additional handling
            if filename.lower().endswith('.ppt'):
                print(f"[DEBUG] Old .ppt format may not be fully supported: {e}")
            else:
                print(f"[DEBUG] Error parsing PowerPoint: {e}")
        
        return chunks

# Định nghĩa hàm standalone cho background task để quản lý vòng đời session riêng
async def run_background_ingestion(doc_id: int, file_path: str, filename: str, version: int, scope: str, target_id: str):
    from app.db.session import AsyncSessionLocal
    import asyncio
    
    print(f"[DEBUG] Background task started for {filename} (ID: {doc_id})")
    # Tạo session mới cho background task
    async with AsyncSessionLocal() as db:
        service = IngestionService(db)
        await service.ingest_file_content(doc_id, file_path, filename, version, scope, target_id)
